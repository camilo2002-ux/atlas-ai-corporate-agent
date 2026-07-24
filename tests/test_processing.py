from __future__ import annotations

import csv
import json
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from atlas.processing.chunking import split_text
from atlas.processing.cleaning import clean_text
from atlas.processing.pipeline import process_document


def _create_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 760, "Política de beneficios de NovaCommerce")
    pdf.drawString(72, 735, "Los colaboradores reciben quince días de vacaciones.")
    pdf.save()


def _create_docx(path: Path) -> None:
    document = Document()
    document.add_heading("Onboarding", level=1)
    document.add_paragraph("El primer día se entrega el equipo corporativo.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Actividad"
    table.cell(0, 1).text = "Plazo"
    table.cell(1, 0).text = "Configurar correo"
    table.cell(1, 1).text = "Día 1"
    document.save(path)


def _create_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Gastos"
    sheet.append(["Categoría", "Límite", "Aprobador"])
    sheet.append(["Hotel", 120, "Gerencia"])
    sheet.append(["Alimentación", 35, "Líder"])
    workbook.save(path)


def _create_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap 2026"
    slide.placeholders[1].text = "Lanzar el asistente Atlas en el tercer trimestre."
    slide.notes_slide.notes_text_frame.text = "Prioridad estratégica de NovaCommerce."
    presentation.save(path)


def _create_md(path: Path) -> None:
    path.write_text(
        "# Incidentes técnicos\n\n## Severidad alta\nEscalar en quince minutos.",
        encoding="utf-8",
    )


def _create_csv(path: Path) -> None:
    with path.open("w", encoding="utf-8", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["ciudad", "horario", "capacidad"])
        writer.writerow(["Quito", "08:00-18:00", "5000"])


def _create_json(path: Path) -> None:
    path.write_text(
        json.dumps(
            {
                "orders_api": {
                    "method": "GET",
                    "path": "/orders/{id}",
                    "auth": "OAuth2",
                }
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )


def _create_html(path: Path) -> None:
    path.write_text(
        "<html><body><h1>Privacidad</h1><p>Los datos se conservan por cinco años.</p></body></html>",
        encoding="utf-8",
    )


@pytest.fixture()
def sample_files(tmp_path: Path) -> list[Path]:
    creators = {
        "sample.pdf": _create_pdf,
        "sample.docx": _create_docx,
        "sample.xlsx": _create_xlsx,
        "sample.pptx": _create_pptx,
        "sample.md": _create_md,
        "sample.csv": _create_csv,
        "sample.json": _create_json,
        "sample.html": _create_html,
    }
    paths: list[Path] = []
    for name, creator in creators.items():
        path = tmp_path / name
        creator(path)
        paths.append(path)
    return paths


def test_clean_text_normalizes_noise() -> None:
    raw = "Título\u00a0 corporativo\r\n\r\n\r\nPágina 1 de 3\r\nContenido   útil"
    assert clean_text(raw) == "Título corporativo\n\nContenido útil"


def test_split_text_creates_overlap() -> None:
    text = " ".join(f"palabra{i}" for i in range(400))
    chunks = split_text(text, max_chars=300, overlap_chars=40)
    assert len(chunks) > 2
    assert all(0 < len(chunk) <= 300 for chunk in chunks)


def test_all_required_formats_produce_chunks(sample_files: list[Path]) -> None:
    for path in sample_files:
        result = process_document(path, max_chars=500, overlap_chars=50)
        assert result.chunks, f"No hubo chunks para {path.suffix}"
        assert all(chunk.metadata["source_file"] == path.name for chunk in result.chunks)
        assert all(chunk.metadata["sha256"] for chunk in result.chunks)


def test_location_metadata_is_preserved(sample_files: list[Path]) -> None:
    expected_location = {
        ".pdf": "page",
        ".docx": "section",
        ".xlsx": "sheet",
        ".pptx": "slide",
        ".md": "section",
        ".csv": "row",
        ".json": "json_path",
        ".html": "section",
    }
    for path in sample_files:
        result = process_document(path)
        key = expected_location[path.suffix]
        assert any(key in chunk.metadata for chunk in result.chunks)
