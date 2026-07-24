from __future__ import annotations

import csv
import json
import re
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from docx import Document
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .cleaning import clean_text
from .models import DocumentUnit

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".html",
    ".htm",
}


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError(f"No fue posible decodificar el archivo: {path.name}")


def _serialize_rows(rows: Iterable[Iterable[Any]]) -> list[str]:
    values = [["" if value is None else str(value) for value in row] for row in rows]
    values = [row for row in values if any(cell.strip() for cell in row)]
    if not values:
        return []

    headers = [cell.strip() or f"columna_{index + 1}" for index, cell in enumerate(values[0])]
    serialized: list[str] = []
    for row_number, row in enumerate(values[1:], start=2):
        padded = row + [""] * max(0, len(headers) - len(row))
        fields = [
            f"{header}: {padded[index].strip()}"
            for index, header in enumerate(headers)
            if index < len(padded) and padded[index].strip()
        ]
        if fields:
            serialized.append(f"Fila {row_number}: " + " | ".join(fields))

    if not serialized:
        serialized.append(" | ".join(headers))
    return serialized


def extract_pdf(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    reader = PdfReader(str(path))
    units: list[DocumentUnit] = []
    warnings: list[str] = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        if len(text) < 30:
            warnings.append(
                f"Página {page_number}: texto mínimo o vacío; puede requerir OCR."
            )
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"page": page_number})
            )

    if not units:
        warnings.append("No se extrajo texto del PDF; probablemente sea escaneado.")
    return units, warnings


def _docx_table_to_text(table: DocxTable) -> str:
    rows = [[cell.text for cell in row.cells] for row in table.rows]
    return "\n".join(_serialize_rows(rows))


def extract_docx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    document = Document(str(path))
    units: list[DocumentUnit] = []
    warnings: list[str] = []
    current_section = "Inicio"
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        text = clean_text("\n\n".join(buffer))
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"section": current_section})
            )
        buffer = []

    if hasattr(document, "iter_inner_content"):
        blocks = document.iter_inner_content()
    else:
        blocks = list(document.paragraphs) + list(document.tables)
        warnings.append(
            "La versión instalada de python-docx no preserva el orden exacto entre párrafos y tablas."
        )

    for block in blocks:
        if isinstance(block, DocxParagraph):
            text = clean_text(block.text)
            if not text:
                continue
            style_name = (block.style.name or "") if block.style else ""
            if style_name.lower().startswith(("heading", "título", "titulo")):
                flush()
                current_section = text
            else:
                buffer.append(text)
        elif isinstance(block, DocxTable):
            table_text = _docx_table_to_text(block)
            if table_text:
                buffer.append(table_text)

    flush()
    return units, warnings


def extract_xlsx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    units: list[DocumentUnit] = []
    warnings: list[str] = []

    try:
        for worksheet in workbook.worksheets:
            rows = list(worksheet.iter_rows(values_only=True))
            serialized = _serialize_rows(rows)
            for index, row_text in enumerate(serialized, start=2):
                units.append(
                    DocumentUnit(
                        text=row_text,
                        metadata={"sheet": worksheet.title, "row": index},
                    )
                )
            if not serialized:
                warnings.append(f"Hoja '{worksheet.title}' sin datos legibles.")
    finally:
        workbook.close()

    return units, warnings


def _pptx_table_text(shape: Any) -> str:
    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
    return "\n".join(_serialize_rows(rows))


def extract_pptx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    presentation = Presentation(str(path))
    units: list[DocumentUnit] = []
    warnings: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        content: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text)
                if text:
                    content.append(text)
            if getattr(shape, "has_table", False):
                table_text = _pptx_table_text(shape)
                if table_text:
                    content.append(table_text)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                notes = clean_text(notes_frame.text)
                if notes:
                    content.append(f"Notas del orador:\n{notes}")

        text = clean_text("\n\n".join(content))
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"slide": slide_number})
            )
        else:
            warnings.append(f"Diapositiva {slide_number} sin texto legible.")

    return units, warnings


def _strip_markdown(text: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"^\s*>\s?", "", text, flags=re.MULTILINE)
    text = re.sub(r"(`{1,3}|\*{1,2}|_{1,2}|~~)", "", text)
    text = re.sub(r"^\s*[-+*]\s+", "", text, flags=re.MULTILINE)
    return clean_text(text)


def extract_markdown(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    raw = _read_text_file(path)
    units: list[DocumentUnit] = []
    current_section = "Inicio"
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        text = _strip_markdown("\n".join(buffer))
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"section": current_section})
            )
        buffer = []

    for line in raw.splitlines():
        heading = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*#*\s*$", line)
        if heading:
            flush()
            current_section = _strip_markdown(heading.group(1)) or "Sin título"
        else:
            buffer.append(line)
    flush()
    return units, []


def extract_csv(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    raw = _read_text_file(path)
    sample = raw[:4096]
    warnings: list[str] = []

    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
        warnings.append("No se detectó el delimitador; se usó coma por defecto.")

    reader = csv.DictReader(raw.splitlines(), dialect=dialect)
    units: list[DocumentUnit] = []
    for row_number, row in enumerate(reader, start=2):
        fields = [
            f"{key}: {value}"
            for key, value in row.items()
            if key and value is not None and str(value).strip()
        ]
        if fields:
            units.append(
                DocumentUnit(
                    text=" | ".join(fields),
                    metadata={"row": row_number},
                )
            )
    return units, warnings


def _flatten_json(value: Any, prefix: str = "$" ) -> list[str]:
    lines: list[str] = []
    if isinstance(value, dict):
        for key, child in value.items():
            lines.extend(_flatten_json(child, f"{prefix}.{key}"))
    elif isinstance(value, list):
        for index, child in enumerate(value):
            lines.extend(_flatten_json(child, f"{prefix}[{index}]"))
    else:
        serialized = json.dumps(value, ensure_ascii=False)
        lines.append(f"{prefix} = {serialized}")
    return lines


def extract_json(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    data = json.loads(_read_text_file(path))
    units: list[DocumentUnit] = []

    if isinstance(data, dict):
        items = [(f"$.{key}", value) for key, value in data.items()]
    elif isinstance(data, list):
        items = [(f"$[{index}]", value) for index, value in enumerate(data)]
    else:
        items = [("$", data)]

    for json_path, value in items:
        text = clean_text("\n".join(_flatten_json(value, json_path)))
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"json_path": json_path})
            )
    return units, []


def _html_table_text(table: Any) -> str:
    rows: list[list[str]] = []
    for row in table.find_all("tr"):
        cells = [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
        if cells:
            rows.append(cells)
    return "\n".join(_serialize_rows(rows))


def extract_html(path: Path) -> tuple[list[DocumentUnit], list[str]]:
    soup = BeautifulSoup(_read_text_file(path), "html.parser")
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.decompose()

    units: list[DocumentUnit] = []
    current_section = "Inicio"
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        text = clean_text("\n\n".join(buffer))
        if text:
            units.append(
                DocumentUnit(text=text, metadata={"section": current_section})
            )
        buffer = []

    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
        if element.name in {"p", "li"} and element.find_parent("table") is not None:
            continue
        if element.name and element.name.startswith("h"):
            flush()
            current_section = element.get_text(" ", strip=True) or "Sin título"
        elif element.name == "table":
            table_text = _html_table_text(element)
            if table_text:
                buffer.append(table_text)
        else:
            text = element.get_text(" ", strip=True)
            if text:
                buffer.append(text)

    flush()
    if not units:
        fallback = clean_text(soup.get_text("\n", strip=True))
        if fallback:
            units.append(DocumentUnit(text=fallback, metadata={"section": "Documento"}))
    return units, []


def extract_file(path: str | Path) -> tuple[list[DocumentUnit], list[str]]:
    source = Path(path)
    suffix = source.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Formato no soportado: {suffix or 'sin extensión'}. "
            f"Permitidos: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    if not source.is_file():
        raise FileNotFoundError(f"No existe el archivo: {source}")

    extractors = {
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".xlsx": extract_xlsx,
        ".pptx": extract_pptx,
        ".md": extract_markdown,
        ".markdown": extract_markdown,
        ".csv": extract_csv,
        ".json": extract_json,
        ".html": extract_html,
        ".htm": extract_html,
    }
    return extractors[suffix](source)
